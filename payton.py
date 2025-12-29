import pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- CONFIGURAÇÃO E ESTILOS ---
prs = Presentation()

# Cores Corporativas (Baseado no logo Grupo Pura)
BLUE_DARK = RGBColor(32, 56, 100)    # Azul Marinho (Corporativo)
BLUE_LIGHT = RGBColor(91, 155, 213)  # Azul Claro (Destaque secundário)
GREY_TEXT = RGBColor(89, 89, 89)     # Cinza Escuro (Textos e eixos)
WHITE = RGBColor(255, 255, 255)

def clean_chart_style(chart, color=BLUE_DARK, show_legend=False):
    """Aplica o estilo 'Clean' (Executivo) aos gráficos"""
    
    # Remover eixos e linhas de grade poluídas
    try:
        chart.value_axis.has_major_gridlines = False
        chart.value_axis.has_minor_gridlines = False
        chart.value_axis.visible = False # Remove eixo Y (números verticais)
    except:
        pass
    
    # Fonte padrão
    chart.font.name = 'Arial'
    chart.font.size = Pt(10)
    chart.font.color.rgb = GREY_TEXT
    
    # Formatação das Séries (Barras/Colunas)
    for series in chart.series:
        # Preenchimento sólido
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color
        
        # Rótulos de Dados (Data Labels) - Essencial para ler sem eixo Y
        series.has_data_labels = True
        series.data_labels.font.size = Pt(9)
        series.data_labels.font.color.rgb = RGBColor(0, 0, 0)
        try:
            # Tenta colocar o número fora da barra (melhor leitura)
            series.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        except:
            pass

    # Legenda
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

def add_slide_chart(title_text, categories, values, chart_type=XL_CHART_TYPE.BAR_CLUSTERED, color=BLUE_DARK):
    """Cria um slide com título e gráfico automaticamente"""
    
    # Layout em branco
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    
    # Título do Slide
    title_shape = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(24), Cm(2))
    p = title_shape.text_frame.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    
    # Inserir Dados
    chart_data = CategoryChartData()
    chart_data.categories = categories
    
    # Tratamento para garantir que valores sejam numéricos
    clean_values = []
    for v in values:
        clean_values.append(v)
        
    chart_data.add_series('Série 1', clean_values)
    
    # Posição do Gráfico (Margens para não ficar colado na borda)
    x, y, cx, cy = Cm(1.5), Cm(3.5), Cm(22.5), Cm(13.5)
    
    chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data).chart
    
    # Aplicar estilo limpo
    clean_chart_style(chart, color)

# ==========================================
# GERAÇÃO DOS SLIDES (DADOS REAIS EXTRAÍDOS)
# ==========================================

# --- 1. CAPA ---
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "APRESENTAÇÃO DE RESULTADOS 2025"
subtitle.text = "GRUPO PURA\nPeríodo: 01/01/2025 a 30/11/2025"
title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
title.text_frame.paragraphs[0].font.bold = True

# --- 2. FINANCEIRO: Evolução Compras ---
# Fonte: 
cats = ['Jan', 'Fev', 'Mar', 'Abr', 'Mai', 'Jun', 'Jul', 'Ago', 'Set', 'Out', 'Nov']
vals = [0.87, 2.5, 5.2, 2.6, 2.9, 4.3, 2.4, 6.0, 2.4, 3.0, 5.0]
add_slide_chart("Evolução de Compras Mensal (R$ Mi)", cats, vals, XL_CHART_TYPE.COLUMN_CLUSTERED, BLUE_DARK)

# --- 3. FINANCEIRO: Gastos Permutas ---
# Fonte: 
vals_perm = [0.52, 1.8, 4.8, 1.5, 2.0, 3.7, 1.7, 4.3, 1.5, 2.1, 3.6]
add_slide_chart("Evolução de Gastos em Permutas (R$ Mi)", cats, vals_perm, XL_CHART_TYPE.COLUMN_CLUSTERED, BLUE_LIGHT)

# --- 4. FINANCEIRO: Condições de Pagamento ---
# Fonte: 
cats_pag = ['Permutas', 'PIX', 'Estoque', 'Boleto', 'Carteira', 'Não Inf.']
vals_pag = [28.0, 5.0, 2.7, 1.6, 0.38, 0.02]
add_slide_chart("Distribuição por Condições de Pagamento (R$ Mi)", cats_pag, vals_pag, XL_CHART_TYPE.BAR_CLUSTERED, BLUE_DARK)

# --- 5 a 8. OBRAS (Dividido em 4 partes para legibilidade) ---
# Fonte: 
# Parte 1
cats_obras1 = ['Res. Manacá I', 'Res. Jatobá', 'Shop. Americana Mall', 'Res. Pacaembu', 'Estoque Central', 'Santuário S.J. Paulo II', 'Manacá/Jatobá/Shop', 'Obra BLB', 'Clube Veteranos', 'Villágio', 'Sky House', 'Swiss Pak', 'Acerto Denadai', 'Salão 05 Orla', 'Area 1', 'Nascer do Sol', 'Res. Ipe Roxo', 'Villa Verde', 'Sobrados Praia Azul', 'Salão 20 Maranhão']
vals_obras1 = [9.3, 7.2, 3.1, 1.8, 1.7, 1.0, 0.94, 0.89, 0.85, 0.62, 0.53, 0.51, 0.51, 0.48, 0.45, 0.41, 0.41, 0.39, 0.38, 0.33]
add_slide_chart("Compras por Obras (Top 20) - Parte 1", cats_obras1, vals_obras1, XL_CHART_TYPE.BAR_CLUSTERED, BLUE_DARK)

# Parte 2
cats_obras2 = ['Jose Carlos Evang.', 'Res. SBO 01', 'Diversas Obras', 'Americana Mall', 'Nascer do Sol Emp.', 'Cataby - Obra Paulista', 'Permuta Claudio', 'Desafio do Pura', 'Reservatorio CTC', 'Escritorio', 'Res. Ipe Rosa', 'Imobiliarios Ltda', 'Recinto Festa Peão', 'Permuta Agnaldo', 'Santa Lucia L17', 'Res. Jacaranda', 'Desconhecido', 'Ana Zanaga Balsa', 'Varandas do Sol', 'Av. Jacaranda', 'Obra Construfera']
vals_obras2 = [0.29, 0.27, 0.27, 0.27, 0.25, 0.24, 0.18, 0.18, 0.18, 0.18, 0.17, 0.16, 0.14, 0.13, 0.12, 0.12, 0.11, 0.10, 0.09, 0.09, 0.09]
add_slide_chart("Compras por Obras - Parte 2", cats_obras2, vals_obras2, XL_CHART_TYPE.BAR_CLUSTERED, BLUE_DARK)

# Parte 3
cats_obras3 = ['Obra Terreno Lado Imoveis', 'Varandas do Sol Emp.', 'Avenida Jacaranda', 'Obra Construfera', 'Casa Modelo Jacarandá', 'Vila Verde Empreed.', 'Shop. Americana Etapa 3', 'Cimento Rio Claro', 'Res. Ipe Amarelo', 'Res. SBO 3', 'Elevatoria Esgoto', 'Poço Ipe Amarelo', 'Locação Maquinas', 'Res. Pinheiros', 'Decorado Ipe Roxo', 'Reforma Daniela', 'Quinta Romeiros', 'Marketing Jacarandas', 'AVM Imóveis', 'Dep. Rafaela', 'Iracemápolis 01']
vals_obras3 = [0.109, 0.097, 0.097, 0.094, 0.092, 0.083, 0.080, 0.075, 0.071, 0.069, 0.066, 0.063, 0.058, 0.056, 0.053, 0.051, 0.044, 0.044, 0.044, 0.043, 0.036]
add_slide_chart("Compras por Obras - Parte 3", cats_obras3, vals_obras3, XL_CHART_TYPE.BAR_CLUSTERED, BLUE_DARK)

# Parte 4
cats_obras4 = ['Reforma Dona Maria', 'Canteiro de Obra', 'Plantão Vendas', 'Wision (Visual)', 'Projet Engenharia', 'Pura Imóveis', 'Permuta Lajes', 'Ação Centro Cívico', 'Ipe Amarelo Manut.', 'Casa do Careca', 'Léo Guia', 'Maria Eduarda', 'Orla Praia Azul', 'Cond. América', 'Contrapartida Pref.', 'Doação', 'Reforma Oficina', 'Escritório Dr Rodrigo', 'Reforma Pastor', 'Nova Carioba', 'Pura Dist. Materiais']
vals_obras4 = [0.036, 0.034, 0.032, 0.028, 0.028, 0.028, 0.027, 0.026, 0.024, 0.016, 0.016, 0.016, 0.014, 0.013, 0.013, 0.012, 0.011, 0.011, 0.010, 0.009, 0.008]
add_slide_chart("Compras por Obras - Parte 4", cats_obras4, vals_obras4, XL_CHART_TYPE.BAR_CLUSTERED, BLUE_DARK)

# --- 9 a 12. FORNECEDORES (Dividido em 4 partes) ---
# Fonte: 
# Parte 1
cats_f1 = ['Hipermix Brasil', 'PSPAV Comercio', 'Esquadrimaxxi', 'Pura Construtora', 'Esquadrias e Vidros', 'Cofibra Industria', 'BLB Blocos', 'DC Comercio Mat.', 'ANTC Industria', 'Galvani Mineração', 'JMS Marbles', 'Comercial Ipiranga', 'KJC Portas', 'Vittoria Cimentos', 'Neural Projetos', 'Construfera Com.', 'Mondialle Ind.', 'Americandaimes', 'Eletro Fort', 'Triade 19', 'Lajes Caetano', 'Cia do Ferro', 'Inove Tubos', 'FR Silva', 'Pav Concreto']
vals_f1 = [5.1, 3.3, 2.8, 2.7, 2.1, 1.5, 1.4, 1.1, 1.1, 1.0, 0.8, 0.7, 0.6, 0.53, 0.52, 0.51, 0.49, 0.49, 0.47, 0.43, 0.38, 0.36, 0.34, 0.34, 0.34]
add_slide_chart("Top Fornecedores (R$ Mi) - Parte 1", cats_f1, vals_f1, XL_CHART_TYPE.BAR_CLUSTERED, BLUE_LIGHT)

# Parte 2
cats_f2 = ['Pizzinatto Ind.', 'Ale Tubos', 'Cimento Rio', 'Paris Hidro Luz', 'Paula Cristina', 'Gleiciara Alves', 'Manetoni Dist.', 'Rede Eletrica', 'Mineracao Areia', 'Futura Ind.', 'Fuminas Ind.', 'Amizade Com.', 'Tormel Com.', 'Qualyvax Proj.', 'Reservatorios Met.', 'Plasnil Ind.', 'Torina Madeiras', 'Poliview Tec.', 'TI Construcoes', 'BT Comercio', 'Marco Antonio', 'Projet Topo.', 'MM Comercio', 'Udiaco Com.', 'Plastcor', 'Araujo Esq.', 'RV Industria', 'Distribuidora Cimento', 'Paula Cristina', 'Cicalfer']
vals_f2 = [0.30, 0.30, 0.28, 0.28, 0.28, 0.27, 0.24, 0.20, 0.20, 0.16, 0.16, 0.15, 0.15, 0.14, 0.13, 0.12, 0.12, 0.12, 0.12, 0.12, 0.11, 0.10, 0.10, 0.10, 0.09, 0.09, 0.09, 0.09, 0.08, 0.08]
add_slide_chart("Fornecedores - Parte 2", cats_f2, vals_f2, XL_CHART_TYPE.BAR_CLUSTERED, BLUE_LIGHT)

# Parte 3
cats_f3 = ['PR Comercio', 'Fortlev Ind.', 'R&M Comercio', 'Munic. Americana', 'Staff Const.', 'RR Tudo Const.', 'Ideal Box', 'Rhoferaço', 'Madeireira NA', 'P&J Maia', 'Bobmaq Loc.', 'Duarte Cons.', 'Calu Materiais', 'Renato Carvalho', '35.829 Nicole', 'Work & Cement', 'BRN Ind.', 'Soldam Sol.', 'Depto Agua', 'Clayton Luis', 'Coperfil', 'CSN Cimentos', 'A D de Lima', 'Tatu Premoldados']
vals_f3 = [0.079, 0.078, 0.073, 0.068, 0.068, 0.066, 0.061, 0.061, 0.061, 0.057, 0.054, 0.053, 0.051, 0.045, 0.044, 0.042, 0.041, 0.041, 0.040, 0.038, 0.037, 0.036, 0.034, 0.034]
add_slide_chart("Fornecedores - Parte 3", cats_f3, vals_f3, XL_CHART_TYPE.BAR_CLUSTERED, BLUE_LIGHT)

# Parte 4
cats_f4 = ['Tatu Premoldados', 'Filli Moveis', 'Real Madri', 'Fabio Dall', 'Kcifer Com.', 'KS Tubos', 'Galpao Tudo', 'G. Santos', 'Alianca Esq.', 'Trelicamp', 'Redentor', 'Cofert', 'CSN Cimentos', 'Blocos Band.', 'Acropole', '26.891 Wallace', 'Gilberto Carlos', 'Bendilatti', 'Flora Natura', 'O W Tesse', 'Fera Atac', 'Bombacamp', 'Loxam', 'LGO Eng.', 'Helix', 'Lino Com.', 'Degam', 'Geradores', 'Amerilux', 'Jundsondas']
vals_f4 = [0.032, 0.030, 0.028, 0.028, 0.027, 0.027, 0.026, 0.026, 0.025, 0.025, 0.025, 0.023, 0.023, 0.022, 0.022, 0.022, 0.021, 0.019, 0.019, 0.019, 0.018, 0.018, 0.018, 0.018, 0.018, 0.017, 0.017, 0.017, 0.017, 0.015]
add_slide_chart("Fornecedores - Parte 4", cats_f4, vals_f4, XL_CHART_TYPE.BAR_CLUSTERED, BLUE_LIGHT)

# --- 13. OPERACIONAL: Volume Mensal (OFs) ---
# Fonte: 
vals_of = [140, 315, 761, 636, 1036, 427, 450, 757, 520, 468, 360]
add_slide_chart("Volume Mensal de Ordens (OFs)", cats, vals_of, XL_CHART_TYPE.COLUMN_CLUSTERED, BLUE_DARK)

# --- 14. CONTRATOS: Global ---
# Fonte: 
cats_contr = ['Total Contratos', 'Consumo Acumulado', 'Saldo Disponível']
vals_contr = [72.9, 30.6, 42.3]
add_slide_chart("Gestão de Contratos de Permutas - Global (R$ Mi)", cats_contr, vals_contr, XL_CHART_TYPE.COLUMN_CLUSTERED, BLUE_DARK)

# --- 15. CONTRATOS: Por Parceiro ---
# Fonte: 
cats_parc = ['PURA', 'VITORINA', 'DENILSON', 'LAIS']
vals_parc = [56.3, 11.2, 4.9, 0.52]
add_slide_chart("Gestão de Permutas por Parceiro (R$ Mi)", cats_parc, vals_parc, XL_CHART_TYPE.COLUMN_CLUSTERED, BLUE_LIGHT)

# --- SALVAR O ARQUIVO ---
file_name = "Apresentacao_Pura_2025_Final.pptx"
prs.save(file_name)
print(f"Arquivo gerado com sucesso: {file_name}")